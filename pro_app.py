import streamlit as st
import hmac
import pandas as pd
import json
import os
import requests
from io import BytesIO # Use BytesIO for pptx
import docx2txt
import PyPDF2
from docx import Document
import time
import re
from datetime import datetime
from pptx import Presentation # Import Presentation for pptx

# --- Authentication ---
def check_password():
    """Returns `True` if the user had a correct password."""

    def login_form():
        """Form with widgets to collect user information"""
        with st.form("Credentials"):
            st.text_input("Username", key="username")
            st.text_input("Password", type="password", key="password")
            st.form_submit_button("Log in", on_click=password_entered)

    def password_entered():
        """Checks whether a password entered by the user is correct."""
        # Ensure secrets and passwords structure exists before accessing
        if "passwords" in st.secrets and st.session_state["username"] in st.secrets["passwords"]:
            stored_password = st.secrets.passwords[st.session_state["username"]]
            # Ensure stored_password is a string or bytes for hmac.compare_digest
            if isinstance(stored_password, (str, bytes)):
                if hmac.compare_digest(
                    st.session_state["password"],
                    str(stored_password) # Ensure it's compared as string if needed
                ):
                    st.session_state["password_correct"] = True
                    del st.session_state["password"]  # Don't store the username or password.
                    del st.session_state["username"]
                    st.session_state["logged_in"] = True # Flag user as logged in
                    return # Exit function on success
                else:
                     st.session_state["password_correct"] = False
            else:
                st.error(f"Password configuration error for user {st.session_state['username']}.")
                st.session_state["password_correct"] = False
        else:
             st.session_state["password_correct"] = False

        # If checks failed or structure doesn't exist
        if not st.session_state.get("password_correct", False):
            st.error("😕 User not known or password incorrect")


    # Initialize login state if not present
    if "logged_in" not in st.session_state:
        st.session_state["logged_in"] = False

    # If user is logged in, show logout button and return True
    if st.session_state.get("logged_in", False):
        if st.sidebar.button("Log out"):
            st.session_state["logged_in"] = False
            st.session_state["password_correct"] = False
            # Optionally clear other session state data on logout if needed
            # st.session_state.documents_content = {}
            # st.session_state.analysis_results = None
            # st.session_state.budget_data = None
            st.rerun() # Rerun to immediately show login form
        return True

    # If not logged in, show login form
    login_form()
    return False

# --- Main Application Logic (only runs if authenticated) ---
if check_password():

    # Set page configuration
    st.set_page_config(
        page_title="Project Health Analysis",
        page_icon="📊",
        layout="wide",
    )

    # Custom CSS (unchanged, CSS properties use American English standard)
    st.markdown("""
    <style>
        .main {
            padding: 2rem;
        }
        .status-green {
            background-color: #d4edda;
            color: #155724;
            padding: 10px;
            border-radius: 5px;
            margin-bottom: 10px;
        }
        .status-amber {
            background-color: #fff3cd;
            color: #856404;
            padding: 10px;
            border-radius: 5px;
            margin-bottom: 10px;
        }
        .status-red {
            background-color: #f8d7da;
            color: #721c24;
            padding: 10px;
            border-radius: 5px;
            margin-bottom: 10px;
        }
        .metric-box {
            background-color: #f8f9fa;
            padding: 15px;
            border-radius: 5px;
            margin-bottom: 15px;
            box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        }
        .score-high {
            color: #155724;
            font-weight: bold;
        }
        .score-medium {
            color: #856404;
            font-weight: bold;
        }
        .score-low {
            color: #721c24;
            font-weight: bold;
        }
    </style>
    """, unsafe_allow_html=True)

    # Initialize session states
    if 'api_key' not in st.session_state:
        st.session_state.api_key = ""
    if 'documents_content' not in st.session_state:
        st.session_state.documents_content = {}
    if 'analysis_results' not in st.session_state:
        st.session_state.analysis_results = None
    if 'budget_data' not in st.session_state:
        st.session_state.budget_data = None

    def extract_text_from_file(file):
        """Extract text content from various file types."""
        file_extension = file.name.split('.')[-1].lower()

        if file_extension == 'txt':
            # Use try-except for decoding robustness
            try:
                return file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                 try:
                    return file.getvalue().decode('latin-1') # Try another common encoding
                 except Exception as e:
                    st.error(f"Error decoding {file.name}: {e}")
                    return None

        elif file_extension == 'docx':
            try:
                # Use BytesIO for docx2txt
                return docx2txt.process(BytesIO(file.getvalue()))
            except Exception as e:
                st.error(f"Error processing DOCX {file.name}: {e}")
                return None

        elif file_extension == 'pdf':
            try:
                pdf_reader = PyPDF2.PdfReader(BytesIO(file.getvalue()))
                text = ""
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n" # Add newline between pages
                return text
            except Exception as e:
                st.error(f"Error processing PDF {file.name}: {e}")
                return None

        elif file_extension == 'pptx':
            try:
                prs = Presentation(BytesIO(file.getvalue()))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                st.error(f"Error processing PPTX {file.name}: {e}")
                return None

        elif file_extension in ['csv', 'xls', 'xlsx']:
            try:
                # Read Excel/CSV into pandas DataFrame
                if file_extension in ['xls', 'xlsx']:
                    df = pd.read_excel(BytesIO(file.getvalue()))
                else: # csv
                     # Attempt to read CSV, handling potential encoding issues
                    try:
                        df = pd.read_csv(BytesIO(file.getvalue()), encoding='utf-8')
                    except UnicodeDecodeError:
                        df = pd.read_csv(BytesIO(file.getvalue()), encoding='latin-1') # Fallback
                return df.to_string()
            except Exception as e:
                st.error(f"Error processing table file {file.name}: {e}")
                return None

        else:
            st.warning(f"Unsupported file type: {file_extension} for file {file.name}")
            return None

    def categorise_document(file_name, file_content):
        """Categorise the document type based on content and filename."""
        file_name_lower = file_name.lower()
        content_lower = file_content.lower() if file_content else "" # Handle None content

        # Simple heuristics for categorisation
        if any(term in file_name_lower for term in ['sow', 'statement', 'work', 'scope']):
            return "Statement of Work"
        elif any(term in file_name_lower for term in ['plan', 'schedule', 'timeline', 'gantt']):
            return "Project Plan"
        elif any(term in file_name_lower for term in ['status', 'report', 'update', 'progress']):
            return "Status Report"
        elif any(term in file_name_lower for term in ['risk', 'issue', 'log', 'raid']):
            return "Risk and Issue Log"
        elif any(term in file_name_lower for term in ['action', 'task', 'todo', 'minutes']):
            return "Action List / Minutes"
        elif any(term in file_name_lower for term in ['budget', 'cost', 'finance', 'financial', 'spend']):
            return "Budget Document"
        elif any(term in file_name_lower for term in ['presentation', 'deck', 'slides']):
             return "Presentation"
        else:
            # Try to determine type from content
            if ('budget' in content_lower or 'financial' in content_lower) and \
               ('$' in content_lower or '€' in content_lower or '£' in content_lower or 'cost' in content_lower):
                return "Budget Document"
            elif 'risk' in content_lower and ('issue' in content_lower or 'mitigation' in content_lower):
                return "Risk and Issue Log"
            elif ('status' in content_lower and 'report' in content_lower) or 'progress update' in content_lower :
                return "Status Report"
            elif ('action item' in content_lower or 'decision log' in content_lower) and \
                 ('assigned' in content_lower or 'due' in content_lower or 'owner' in content_lower):
                return "Action List / Minutes"
            elif 'scope' in content_lower and ('deliverable' in content_lower or 'objective' in content_lower or 'requirement' in content_lower):
                return "Statement of Work"
            elif 'timeline' in content_lower or 'milestone' in content_lower or 'gantt' in content_lower:
                 return "Project Plan"
            else:
                return "Other Document"

    def extract_budget_info(documents_content):
        """Extract budget information from documents."""
        budget_data = {
            "total_budget": None,
            "spent": None,
            "remaining": None,
            "over_under": None,
            "details": []
        }

        # Refined regex patterns to capture currency symbols optionally at the start or end
        budget_pattern = r'(?:budget|total cost|project cost|allocated budget|approved budget)[\s:]*(?:[$£€]\s?)?([\d,]+(?:\.\d{1,2})?)\s?(?:[$£€])?'
        spent_pattern = r'(?:spent|expenses|costs to date|actual cost|expenditure)[\s:]*(?:[$£€]\s?)?([\d,]+(?:\.\d{1,2})?)\s?(?:[$£€])?'

        # First pass to find budget documents
        budget_docs = {name: content for name, content in documents_content.items()
                       if categorise_document(name, content) == "Budget Document"}

        # If no explicit budget documents, search all documents
        search_docs = budget_docs if budget_docs else documents_content

        found_budgets = []
        found_spent = []

        for doc_name, content in search_docs.items():
            if not content: continue # Skip if content is None

            # Look for budget figures
            # Use re.IGNORECASE for case-insensitivity
            budget_matches = re.findall(budget_pattern, content, re.IGNORECASE)
            spent_matches = re.findall(spent_pattern, content, re.IGNORECASE)

            if budget_matches:
                for match in budget_matches:
                    try:
                        budget_value = float(str(match).replace(',', ''))
                        found_budgets.append(budget_value)
                        budget_data["details"].append({
                            "document": doc_name,
                            "type": "budget",
                            "value_found": budget_value
                        })
                    except (ValueError, TypeError):
                        pass # Ignore if conversion fails

            if spent_matches:
                 for match in spent_matches:
                    try:
                        spent_value = float(str(match).replace(',', ''))
                        found_spent.append(spent_value)
                        budget_data["details"].append({
                            "document": doc_name,
                            "type": "spent",
                            "value_found": spent_value
                        })
                    except (ValueError, TypeError):
                        pass # Ignore if conversion fails

        # Determine the most likely total budget and amount spent (often the highest value mentioned)
        if found_budgets:
             budget_data["total_budget"] = max(found_budgets)
        if found_spent:
             budget_data["spent"] = max(found_spent)


        # Calculate remaining and over/under
        if budget_data["total_budget"] is not None and budget_data["spent"] is not None:
            # Ensure values are reasonable (e.g., spent isn't vastly larger than budget unless it's explicitly stated as overspend)
            # This logic might need refinement based on typical document structures
            budget_data["remaining"] = budget_data["total_budget"] - budget_data["spent"]
            budget_data["over_under"] = "under" if budget_data["remaining"] >= 0 else "over"

        return budget_data


    def call_claude_api(api_key, prompt, model="claude-3-sonnet-20240229"):
        """Call the Claude API directly using requests"""
        url = "https://api.anthropic.com/v1/messages"

        headers = {
            "Content-Type": "application/json",
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01" # Use a recommended stable version
        }

        data = {
            "model": model,
            "max_tokens": 4000, # Max output tokens
            "temperature": 0.1, # Lower temperature for more factual responses
            "system": "You are a project management expert analysing project documentation. Provide clear, objective analysis based only on the facts presented in the documents. Use British English spelling (e.g., analyse, categorise).",
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=180) # Increased timeout
            response.raise_for_status() # Raise HTTPError for bad responses (4xx or 5xx)

            response_json = response.json()

            if 'content' in response_json and len(response_json['content']) > 0 and 'text' in response_json['content'][0]:
                return response_json['content'][0]['text']
            elif 'error' in response_json:
                 st.error(f"API Error: {response_json['error'].get('type')} - {response_json['error'].get('message')}")
                 return None
            else:
                st.error("Unexpected API response format.")
                st.json(response_json) # Show the unexpected response
                return None

        except requests.exceptions.RequestException as e:
            st.error(f"Error calling Claude API: {e}")
            # Log more details if possible, e.g., e.response.text
            if hasattr(e, 'response') and e.response is not None:
                 st.error(f"Response status: {e.response.status_code}")
                 st.error(f"Response text: {e.response.text}")
            return None
        except Exception as e:
            st.error(f"An unexpected error occurred during API call: {str(e)}")
            return None

    def analyse_project_with_claude(api_key, documents_content):
        """Send project documents to Claude for analysis."""
        # Define budget_data initially to avoid UnboundLocalError
        budget_data = extract_budget_info(documents_content)

        try:
            # Prepare documents for Claude
            docs_formatted = ""
            total_chars = 0
            char_limit = 150000 # Adjust based on typical token limits (~4 chars/token)

            for filename, content in documents_content.items():
                 if not content: continue # Skip empty content

                 doc_type = categorise_document(filename, content)
                 # Truncate content smartly if needed
                 content_to_add = f"\n\n--- DOCUMENT: {filename} (Type: {doc_type}) ---\n{content}"
                 if total_chars + len(content_to_add) > char_limit:
                     remaining_chars = char_limit - total_chars
                     if remaining_chars > 200: # Add truncated content if space allows
                         content_to_add = content_to_add[:remaining_chars] + "... [TRUNCATED]"
                     else:
                         st.warning(f"Content limit reached, skipping document: {filename}")
                         continue # Skip this document if too close to limit

                 docs_formatted += content_to_add
                 total_chars += len(content_to_add)


            budget_info = "\nBudget Summary (automatically extracted):\n"
            if budget_data.get("total_budget") is not None:
                budget_info += f"- Total Budget: £{budget_data['total_budget']:,.2f}\n" # Assuming GBP, adjust if needed
                if budget_data.get("spent") is not None:
                    budget_info += f"- Spent: £{budget_data['spent']:,.2f}\n"
                    budget_info += f"- Remaining: £{budget_data['remaining']:,.2f}\n"
                    budget_info += f"- Status: {budget_data['over_under']}spend\n"
                else:
                    budget_info += "- Spent: Not clearly identified\n"
            else:
                budget_info += "- Total Budget: Not clearly identified\n"
            budget_info += f"- Budget Details Found: {len(budget_data.get('details',[]))} potential figures identified across documents."


            # Prompt for Claude (using British English)
            prompt = f"""
            You are a project management expert reviewing project documentation. Analyse the provided project documents looking for these key aspects:

            1. Scope creep indicators (unplanned changes, new requests not in original scope)
            2. Dependency mapping quality (are dependencies clearly identified and tracked?)
            3. Objective and goal setting quality (are objectives SMART - Specific, Measurable, Achievable, Relevant, Time-bound?)
            4. Budget situation (based on provided summary and document content)
            5. Planning quality (is the plan detailed, realistic, with clear milestones?)
            6. Key risks and issues (identify major threats and problems)

            Based on your analysis, provide the following outputs:

            1. Scope Creep: List specific instances or indicators of potential scope creep (as bullet points). If none, state "No specific scope creep indicators identified".
            2. Dependency Mapping: Score the quality from 1-10 (1=poor, 10=excellent) and briefly explain your reasoning.
            3. Objective Setting: Score the quality from 1-10 and provide specific examples of good/poor objectives found in the documents.
            4. Budget Analysis: Briefly analyse the budget situation based on the summary provided ({budget_info}) and any further details in the documents. Mention any discrepancies or concerns.
            5. Planning Quality: Score the quality from 1-10 and explain your reasoning.
            6. Risks & Issues: Identify the top 3-5 most significant risks and issues mentioned.
            7. Project Status: Based *only* on the information in the documents, classify the project's overall health as:
               - GREEN (Appears on track, minor issues)
               - AMBER (Some concerns, potential risks need attention)
               - RED (Significant issues, likely off-track)
               Provide a concise justification (1-2 sentences) for this status, referencing specific evidence (e.g., missed deadlines, budget overspend, critical risks).

            Format your response *strictly* as a JSON object with these keys:
            "scope_creep_items": ["item1", "item2", ...],
            "dependency_mapping_score": number,
            "dependency_mapping_reasoning": "string",
            "objective_setting_score": number,
            "objective_setting_reasoning": "string",
            "objective_examples": {{"good": ["example1", ...], "poor": ["example1", ...]}},
            "budget_analysis": "string",
            "planning_quality_score": number,
            "planning_quality_reasoning": "string",
            "top_risks_issues": ["item1", "item2", ...],
            "project_status": "GREEN" | "AMBER" | "RED",
            "status_justification": "string"

            Ensure the entire output is a single, valid JSON object starting with {{ and ending with }}. Use British English spellings.

            Documents to analyse:
            {docs_formatted}
            """

            # Get model from secrets if available
            model = "claude-3-sonnet-20240229" # Default model
            try:
                if "ANTHROPIC_MODEL" in st.secrets and st.secrets["ANTHROPIC_MODEL"]:
                    model = st.secrets["ANTHROPIC_MODEL"]
            except Exception as e:
                 st.warning(f"Could not read ANTHROPIC_MODEL from secrets: {e}. Using default.")


            # Make a direct API call to Claude
            response_text = call_claude_api(api_key, prompt, model)

            if response_text:
                # Attempt to find and parse JSON object in the response
                # Handle potential markdown ```json ... ``` blocks
                if response_text.strip().startswith("```json"):
                     json_str = response_text.strip()[7:-3].strip() # Remove markdown fences
                elif response_text.strip().startswith("{") and response_text.strip().endswith("}"):
                    json_str = response_text.strip()
                else:
                    # Try finding JSON within the text if not perfectly formatted
                    json_start = response_text.find('{')
                    json_end = response_text.rfind('}') + 1
                    if json_start != -1 and json_end > json_start:
                         json_str = response_text[json_start:json_end]
                    else:
                         st.error("Claude did not return a recognisable JSON response.")
                         st.text_area("Raw API Response:", response_text, height=200)
                         return None, budget_data # Return budget data even if analysis fails


                try:
                    analysis_results = json.loads(json_str)
                    # Basic validation of expected keys
                    expected_keys = ["scope_creep_items", "dependency_mapping_score", "objective_setting_score",
                                     "budget_analysis", "planning_quality_score", "top_risks_issues",
                                     "project_status", "status_justification"]
                    if all(key in analysis_results for key in expected_keys):
                         return analysis_results, budget_data
                    else:
                         st.error("Parsed JSON is missing some expected keys.")
                         st.json(analysis_results) # Show what was parsed
                         st.text_area("Original String Parsed:", json_str, height=100)
                         return None, budget_data

                except json.JSONDecodeError as e:
                    st.error(f"Failed to parse Claude's JSON response: {e}")
                    st.text_area("Problematic JSON String:", json_str, height=200)
                    st.text_area("Full Raw API Response:", response_text, height=100)
                    return None, budget_data
            else:
                # Error message already shown in call_claude_api
                # st.error("Failed to get a valid response from Claude API") # Redundant
                return None, budget_data

        except Exception as e:
            st.error(f"Error in project analysis process: {str(e)}")
            import traceback
            st.error(traceback.format_exc()) # Print stack trace for debugging
            return None, budget_data # Return budget data even if analysis fails

    def display_results(analysis_results, budget_data):
        """Display the analysis results in a formatted Streamlit interface."""
        if not analysis_results:
            st.error("No analysis results to display.")
            return

        # Project Status Header with appropriate styling
        status = analysis_results.get("project_status", "UNKNOWN").upper()
        status_justification = analysis_results.get('status_justification', 'No justification provided.')

        if status == "GREEN":
            st.markdown(f"""<div class="status-green">
                          <h2>Project Status: 🟢 GREEN</h2>
                          <p>{status_justification}</p>
                        </div>""", unsafe_allow_html=True)
        elif status == "AMBER":
            st.markdown(f"""<div class="status-amber">
                          <h2>Project Status: 🟡 AMBER</h2>
                           <p>{status_justification}</p>
                        </div>""", unsafe_allow_html=True)
        elif status == "RED":
            st.markdown(f"""<div class="status-red">
                          <h2>Project Status: 🔴 RED</h2>
                           <p>{status_justification}</p>
                        </div>""", unsafe_allow_html=True)
        else:
             st.warning(f"Unknown Project Status: {status}")
             st.markdown(f"<p>{status_justification}</p>", unsafe_allow_html=True)


        # Create two columns for layout
        col1, col2 = st.columns(2)

        # Column 1: Scores and Objectives/Planning
        with col1:
            st.subheader("Quality Scores (out of 10)")
            # Dependency Mapping Score
            dep_score = analysis_results.get("dependency_mapping_score", "N/A")
            st.markdown(f"""<div class="metric-box">
                          <h4>Dependency Mapping</h4>
                          <p class="{get_score_class(dep_score)}">Score: {dep_score}/10</p>
                          <p><i>Reasoning:</i> {analysis_results.get('dependency_mapping_reasoning', 'N/A')}</p>
                        </div>""", unsafe_allow_html=True)

            # Objective Setting Score
            obj_score = analysis_results.get("objective_setting_score", "N/A")
            st.markdown(f"""<div class="metric-box">
                          <h4>Objective Setting (SMART)</h4>
                           <p class="{get_score_class(obj_score)}">Score: {obj_score}/10</p>
                          <p><i>Reasoning:</i> {analysis_results.get('objective_setting_reasoning', 'N/A')}</p>
                        </div>""", unsafe_allow_html=True)

            # Planning Quality Score
            plan_score = analysis_results.get("planning_quality_score", "N/A")
            st.markdown(f"""<div class="metric-box">
                          <h4>Planning Quality</h4>
                           <p class="{get_score_class(plan_score)}">Score: {plan_score}/10</p>
                           <p><i>Reasoning:</i> {analysis_results.get('planning_quality_reasoning', 'N/A')}</p>
                        </div>""", unsafe_allow_html=True)

            # Objective Examples
            st.subheader("Objective Examples")
            obj_examples = analysis_results.get('objective_examples', {})
            good_objectives = obj_examples.get('good', [])
            poor_objectives = obj_examples.get('poor', [])

            expander_obj = st.expander("View Objective Examples")
            with expander_obj:
                 st.markdown("##### Good Examples:")
                 if good_objectives:
                     for example in good_objectives:
                         st.markdown(f"- {example}")
                 else:
                     st.markdown("_No specific good examples identified._")

                 st.markdown("##### Poor Examples:")
                 if poor_objectives:
                     for example in poor_objectives:
                         st.markdown(f"- {example}")
                 else:
                     st.markdown("_No specific poor examples identified._")


        # Column 2: Scope Creep, Risks, Budget
        with col2:
            # Scope Creep Items
            st.subheader("Scope Creep Indicators")
            scope_creep_items = analysis_results.get("scope_creep_items", [])
            if scope_creep_items and scope_creep_items[0] != "No specific scope creep indicators identified":
                st.markdown(f"""<div class="metric-box">""", unsafe_allow_html=True)
                for item in scope_creep_items:
                    st.markdown(f"- {item}")
                st.markdown("</div>", unsafe_allow_html=True)
            else:
                st.markdown("""<div class="metric-box">
                                <p><i>No specific scope creep indicators identified.</i></p>
                             </div>""", unsafe_allow_html=True)


            # Top Risks and Issues
            st.subheader("Top Risks & Issues")
            top_risks = analysis_results.get("top_risks_issues", [])
            if top_risks:
                st.markdown(f"""<div class="metric-box">""", unsafe_allow_html=True)
                for i, item in enumerate(top_risks, 1):
                    st.markdown(f"{i}. {item}")
                st.markdown("</div>", unsafe_allow_html=True)
            else:
                 st.markdown("""<div class="metric-box">
                                <p><i>No specific risks or issues highlighted by the analysis.</i></p>
                             </div>""", unsafe_allow_html=True)


            # Budget Analysis
            st.subheader("Budget Analysis")
            st.markdown("""<div class="metric-box">""", unsafe_allow_html=True)
            budget_analysis_text = analysis_results.get('budget_analysis', 'Budget analysis not provided.')
            st.markdown(f"<p>{budget_analysis_text}</p>", unsafe_allow_html=True)

            # Display extracted budget figures if available
            if budget_data and budget_data.get("total_budget") is not None:
                st.markdown("---")
                st.markdown("##### Extracted Budget Summary:")
                st.markdown(f"<p>Total Budget: <strong>£{budget_data['total_budget']:,.2f}</strong></p>", unsafe_allow_html=True) # Assuming GBP

                if budget_data.get("spent") is not None:
                    over_under_class = "score-high" if budget_data["over_under"] == "under" else "score-low"
                    st.markdown(f"""
                        <p>Spent to Date: <strong>£{budget_data['spent']:,.2f}</strong></p>
                        <p>Remaining: <strong>£{budget_data['remaining']:,.2f}</strong></p>
                        <p class="{over_under_class}">Status: {budget_data['over_under'].upper()}SPEND</p>
                    """, unsafe_allow_html=True)
                else:
                     st.markdown("<p>Spent to Date: <i>Not clearly identified</i></p>", unsafe_allow_html=True)
            elif budget_data and budget_data.get("details"):
                 st.markdown("---")
                 st.markdown("##### Extracted Budget Summary:")
                 st.markdown("<p><i>Total budget or amount spent could not be definitively determined, but potential figures were found.</i></p>", unsafe_allow_html=True)
                 # Optionally show details found
                 # with st.expander("Show details"):
                 #    st.json(budget_data["details"])

            st.markdown("</div>", unsafe_allow_html=True) # Close metric-box


    def get_score_class(score):
        """Return CSS class based on score value."""
        try:
            score_num = float(score)
            if score_num >= 7:
                return "score-high"
            elif score_num >= 4:
                return "score-medium"
            else:
                return "score-low"
        except (ValueError, TypeError):
             return "" # No specific class if score is not a number


    # --- UI Rendering ---

    st.title("Project Health Analysis")
    st.markdown("Upload project documents to analyse key project management indicators using AI.")

    # Sidebar configuration - Handle API key from secrets or input
    with st.sidebar:
        st.header("Configuration")

        # API Key Input (only if not found in secrets)
        api_key_source = "secrets"
        try:
            if "ANTHROPIC_API_KEY" in st.secrets and st.secrets["ANTHROPIC_API_KEY"]:
                st.session_state.api_key = st.secrets["ANTHROPIC_API_KEY"]
                # Don't display the key, just confirm it's loaded
                st.success("✅ Anthropic API key loaded securely.")
            else:
                 api_key_source = "input"
                 st.session_state.api_key = st.text_input("Anthropic API Key", value=st.session_state.get('api_key',""), type="password", help="Required for analysis.")
        except Exception as e:
            st.warning(f"Could not check secrets for API key: {e}. Please enter it below.")
            api_key_source = "input"
            st.session_state.api_key = st.text_input("Anthropic API Key", value=st.session_state.get('api_key',""), type="password", help="Required for analysis.")


        st.markdown("---")
        st.markdown("### Supported Document Types")
        st.markdown("""
        - Text (`.txt`)
        - PDF (`.pdf`)
        - Word (`.docx`)
        - PowerPoint (`.pptx`)
        - Excel (`.xlsx`, `.xls`)
        - CSV (`.csv`)
        """)
        st.markdown("The tool attempts to automatically categorise uploaded documents.")

        st.markdown("---")
        st.markdown("### About")
        st.markdown("""
        This app uses Claude AI to analyse project documents for insights on:
        - Scope creep
        - Dependency mapping
        - Objective quality (SMART)
        - Budget situation
        - Planning quality
        - Key risks and issues

        An overall project health status (RED/AMBER/GREEN) is provided based on the analysis. Remember that AI analysis is a tool and should be reviewed by a human expert.
        """)


    # --- Main Page Content ---

    # Upload Section
    st.header("1. Document Upload")
    uploaded_files = st.file_uploader(
        "Upload Project Documents",
        accept_multiple_files=True,
        # Add pptx to supported types
        type=["txt", "pdf", "docx", "csv", "xlsx", "xls", "pptx"],
        help="Upload relevant project files like SoWs, Plans, Status Reports, Risk Logs, Budgets, Presentations etc."
    )

    # File Processing and Display Area
    if uploaded_files:
        docs_processed = {}
        processing_errors = False
        st.markdown("---")
        st.markdown("#### Processing Uploaded Files...")
        progress_bar = st.progress(0)
        status_text = st.empty()

        for i, uploaded_file in enumerate(uploaded_files):
            status_text.text(f"Processing {uploaded_file.name}...")
            content = extract_text_from_file(uploaded_file)
            if content is not None: # Check if extraction was successful (not None)
                # Only add successfully processed files to session state
                doc_type = categorise_document(uploaded_file.name, content)
                st.session_state.documents_content[uploaded_file.name] = content
                docs_processed[uploaded_file.name] = doc_type
                status_text.text(f"Processed {uploaded_file.name} (Detected as: {doc_type})")
            else:
                 # Error message handled within extract_text_from_file
                 status_text.text(f"Failed to process {uploaded_file.name}. Check error messages above.")
                 processing_errors = True
            progress_bar.progress((i + 1) / len(uploaded_files))

        status_text.text("File processing complete.")
        st.markdown("#### Successfully Processed Documents:")
        if docs_processed:
             # Create a simple table or list
             df_processed = pd.DataFrame(docs_processed.items(), columns=['Filename', 'Detected Type'])
             st.dataframe(df_processed, use_container_width=True)
        else:
             st.info("No documents were successfully processed.")
        if processing_errors:
            st.warning("Some files could not be processed. Analysis will run on successfully processed files only.")


    # Analysis Trigger Section
    st.markdown("---")
    st.header("2. Run Analysis")

    # Check if files are ready for analysis
    if st.session_state.documents_content:
        col_run, col_clear = st.columns(2)

        with col_run:
             # Disable button if API key is missing
             disable_analysis = not st.session_state.api_key
             if disable_analysis and api_key_source == "input":
                 st.warning("Please enter your Anthropic API Key in the sidebar to enable analysis.")

             analyze_button = st.button("Analyse Project Health", disabled=disable_analysis, type="primary", help="Sends document content to Claude AI for analysis.")

        with col_clear:
            # Button to clear uploaded files and previous results
            if st.button("Clear Uploaded Files & Results", help="Removes uploaded documents and any previous analysis results."):
                st.session_state.documents_content = {}
                st.session_state.analysis_results = None
                st.session_state.budget_data = None
                st.success("Cleared uploaded documents and results.")
                time.sleep(1) # Brief pause before rerun
                st.rerun()


        if analyze_button and not disable_analysis:
            with st.spinner("Analysing project documentation with AI... This might take a minute or two depending on document size."):
                analysis_results, budget_data = analyse_project_with_claude(
                    st.session_state.api_key,
                    st.session_state.documents_content
                )

                if analysis_results:
                    st.session_state.analysis_results = analysis_results
                    st.session_state.budget_data = budget_data # Store budget data alongside results
                    st.success("Analysis complete!")
                    # Rerun to display results immediately without waiting for another interaction
                    st.rerun()
                else:
                    # Errors handled within the analysis function
                    st.error("Analysis failed. Please check the error messages above and ensure your API key is correct.")
                    # Keep any existing analysis results if the new run fails
                    # st.session_state.analysis_results = None # Avoid clearing potentially useful old results
                    # st.session_state.budget_data = None


    elif uploaded_files:
         st.warning("No documents were successfully processed for analysis. Please check file types and content.")
    else:
        st.info("👆 Upload project documents to begin.")


    # Display Results Section
    if st.session_state.analysis_results:
        st.markdown("---")
        st.header("3. Project Health Report")
        display_results(st.session_state.analysis_results, st.session_state.budget_data)

        # Add Clear Analysis Results Button specifically for the report section
        if st.button("Clear Analysis Report", help="Clears only the displayed report, keeping uploaded files."):
                st.session_state.analysis_results = None
                st.session_state.budget_data = None
                st.success("Analysis report cleared.")
                time.sleep(1)
                st.rerun()


        # Add download button for report
        st.markdown("---")
        st.subheader("Download Report")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create a formatted report for download (using British English where appropriate)
        # (Keep existing report generation logic, ensuring consistency with displayed results)
        try: # Wrap report generation in try-except
            report_md = f"""# Project Health Analysis Report
Generated on: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
Documents Analysed: {len(st.session_state.get('documents_content', {}))}

## Overall Project Status: {st.session_state.analysis_results.get('project_status', 'UNKNOWN')}
**Justification:** {st.session_state.analysis_results.get('status_justification', 'N/A')}

## Quality Scores (out of 10)
- **Dependency Mapping:** {st.session_state.analysis_results.get('dependency_mapping_score', 'N/A')}/10
  - *Reasoning:* {st.session_state.analysis_results.get('dependency_mapping_reasoning', 'N/A')}
- **Objective Setting (SMART):** {st.session_state.analysis_results.get('objective_setting_score', 'N/A')}/10
  - *Reasoning:* {st.session_state.analysis_results.get('objective_setting_reasoning', 'N/A')}
- **Planning Quality:** {st.session_state.analysis_results.get('planning_quality_score', 'N/A')}/10
  - *Reasoning:* {st.session_state.analysis_results.get('planning_quality_reasoning', 'N/A')}

## Scope Creep Indicators
"""
            scope_items = st.session_state.analysis_results.get('scope_creep_items', [])
            if scope_items and scope_items[0] != "No specific scope creep indicators identified":
                for item in scope_items:
                    report_md += f"- {item}\n"
            else:
                report_md += "- *No specific scope creep indicators identified.*\n"

            report_md += "\n## Top Risks and Issues\n"
            risk_items = st.session_state.analysis_results.get('top_risks_issues', [])
            if risk_items:
                 for i, item in enumerate(risk_items, 1):
                    report_md += f"{i}. {item}\n"
            else:
                 report_md += "- *No specific risks or issues highlighted by the analysis.*\n"


            report_md += f"\n## Budget Analysis\n"
            report_md += f"{st.session_state.analysis_results.get('budget_analysis', 'Budget analysis not provided.')}\n"
            # Add extracted summary if available
            if st.session_state.budget_data and st.session_state.budget_data.get("total_budget") is not None:
                report_md += f"\n**Extracted Summary:**\n"
                report_md += f"- Total Budget: £{st.session_state.budget_data['total_budget']:,.2f}\n" # Assuming GBP
                if st.session_state.budget_data.get("spent") is not None:
                    report_md += f"- Spent to Date: £{st.session_state.budget_data['spent']:,.2f}\n"
                    report_md += f"- Remaining: £{st.session_state.budget_data['remaining']:,.2f}\n"
                    report_md += f"- Status: {st.session_state.budget_data['over_under'].upper()}SPEND\n"
                else:
                    report_md += "- Spent to Date: Not clearly identified\n"


            report_md += f"\n## Objective Setting Examples\n"
            obj_examples_dl = st.session_state.analysis_results.get('objective_examples', {})
            good_obj_dl = obj_examples_dl.get('good', [])
            poor_obj_dl = obj_examples_dl.get('poor', [])
            report_md += "### Good Examples:\n"
            if good_obj_dl:
                 for ex in good_obj_dl: report_md += f"- {ex}\n"
            else: report_md += "- *None identified*\n"
            report_md += "\n### Poor Examples:\n"
            if poor_obj_dl:
                 for ex in poor_obj_dl: report_md += f"- {ex}\n"
            else: report_md += "- *None identified*\n"

            st.download_button(
                label="Download Report (Markdown)",
                data=report_md.encode('utf-8'), # Encode to bytes
                file_name=f"project_health_report_{timestamp}.md",
                mime="text/markdown"
            )
        except Exception as e:
             st.error(f"Error generating download report: {e}")


# --- Footer or Instructions if nothing else is shown ---
# (Login form is shown by check_password() if not logged in)
# (Instructions to upload are shown if logged in but no files/results)
